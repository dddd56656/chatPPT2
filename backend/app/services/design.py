"""
设计美化服务模块 (design.py) - Google Standard Refactor
"""

import os
import logging
import io
from typing import Dict, Any, List, Optional, Union
from pptx import Presentation
from pptx.slide import SlideLayout, Slide
from .exporter import PPTExporter

logger = logging.getLogger(__name__)

class TemplateEngine:
    def __init__(self, exporter: PPTExporter, template_path: str = "templates/business_report.pptx"):
        self.exporter = exporter
        self.template_path = template_path

    def create_from_template(self, title: str, slides_data: List[Dict[str, Any]]) -> Dict[str, Any]:
        """安全地应用模板，包含降级策略"""
        prs = self._load_resource()
        self._clear_existing_slides(prs)

        # 1. 至少创建一个标题页
        self._create_slide(prs, {"slide_type": "title", "title": title, "subtitle": "AI Generated Presentation"})

        # 2. 遍历生成内容
        if slides_data:
            for i, slide_info in enumerate(slides_data):
                try:
                    self._create_slide(prs, slide_info)
                except Exception as e:
                    logger.error(f"Slide {i+1} generation failed, skipping: {e}")
                    # 容错：创建一个只有标题的空白页，避免整个文件失败
                    self._create_simple_slide(prs, f"Error: Slide {i+1}", str(e))

        return self.exporter.export_ppt(prs, title=title)

    def _load_resource(self) -> Presentation:
        if os.path.exists(self.template_path):
            return Presentation(self.template_path)
        logger.warning(f"Template not found at {self.template_path}, using default.")
        return Presentation()

    def _clear_existing_slides(self, prs: Presentation):
        """清空初始模板中的幻灯片"""
        xml_slides = prs.slides._sldIdLst
        slides_count = len(xml_slides)
        for i in range(slides_count - 1, -1, -1):
            rId = xml_slides[i].rId
            prs.part.drop_rel(rId)
            del xml_slides[i]

    def _create_slide(self, prs: Presentation, info: Dict[str, Any]):
        """统一的幻灯片工厂方法"""
        s_type = info.get("slide_type", "content")
        title = info.get("title", "Untitled")

        # 策略模式：根据类型选择处理逻辑
        if s_type == "title":
            layout = self._find_layout(prs, ["Title Slide", "Title", "Title 1"], 0)
            slide = prs.slides.add_slide(layout)
            self._set_text(slide.shapes.title, title)
            # 尝试找副标题框
            if len(slide.placeholders) > 1:
                self._set_text(slide.placeholders[1], info.get("subtitle", ""))

        elif s_type == "two_column":
            # 尝试找两栏布局，找不到则降级为普通内容页
            layout = self._find_layout(prs, ["Two Content", "Two Column", "Comparison"], None)
            if layout:
                slide = prs.slides.add_slide(layout)
                self._set_text(slide.shapes.title, title)
                placeholders = self._get_body_placeholders(slide)
                if len(placeholders) >= 2:
                    self._set_text(placeholders[0], info.get("left_content", []))
                    self._set_text(placeholders[1], info.get("right_content", []))
            else:
                # 降级策略
                logger.warning(f"Two-column layout missing for '{title}', downgrading to content.")
                # 合并内容
                merged_content = info.get("left_content", []) + info.get("right_content", [])
                self._create_slide(prs, {"slide_type": "content", "title": title, "content": merged_content})

        else: # default content
            layout = self._find_layout(prs, ["Title and Content", "Content"], 1)
            slide = prs.slides.add_slide(layout)
            self._set_text(slide.shapes.title, title)
            placeholders = self._get_body_placeholders(slide)
            if placeholders:
                self._set_text(placeholders[0], info.get("content", []))

    def _create_simple_slide(self, prs: Presentation, title: str, body: str):
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = body

    def _find_layout(self, prs: Presentation, names: List[str], fallback_index: Optional[int]) -> SlideLayout:
        """按名称模糊匹配布局"""
        for layout in prs.slide_layouts:
            for name in names:
                if name.lower() in layout.name.lower():
                    return layout
        
        if fallback_index is not None and fallback_index < len(prs.slide_layouts):
            return prs.slide_layouts[fallback_index]
        return prs.slide_layouts[0] # 最终保底

    def _get_body_placeholders(self, slide: Slide) -> List[Any]:
        """获取非标题的占位符"""
        return [p for p in slide.placeholders if p.element.ph_idx > 0] # 0通常是标题

    def _set_text(self, shape, content):
        """智能填充文本或列表"""
        if not shape or not hasattr(shape, "text_frame"): return
        tf = shape.text_frame
        tf.clear()
        
        if isinstance(content, list):
            for item in content:
                p = tf.add_paragraph()
                p.text = str(item)
                p.level = 0
        else:
            tf.text = str(content)
