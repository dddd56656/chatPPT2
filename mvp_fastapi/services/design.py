"""
设计美化服务模块。

该模块负责应用模板和样式到PPT内容，支持Office官方免费模板。
Refactored according to Google Engineering Practices.
"""

import os
import logging
import io
from typing import Dict, Any, List, Optional, Union

from pptx import Presentation
from pptx.slide import SlideLayout, Slide
from pptx.util import Inches


# 假设 exporter 模块存在 (使用 Mock 来避免循环依赖)
# 在实际运行时，应保证 'from .exporter import PPTExporter' 可用
class PPTExporter:
    """Mock Exporter for testing purposes."""

    def export_ppt(self, prs: Presentation, title: str) -> Dict[str, Any]:
        # 模拟导出：将 prs 保存到 BytesIO
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return {"buffer": buffer, "filename": f"{title}.pptx"}


# 配置日志
logger = logging.getLogger(__name__)


class TemplateEngine:
    """PPT模板引擎类，负责应用模板和样式到PPT内容。"""

    # 根据用户提供的日志，适配模板布局名称
    LAYOUT_NAMES = {
        "title": "Title 1",  # LAYOUT 0
        "content": "Title and Content",  # LAYOUT 4
        "two_column": "Title and two Content 1",  # LAYOUT 6
    }

    # 回退索引，防止名称匹配失败
    LAYOUT_INDICES = {"title": 0, "content": 4, "two_column": 6}

    def __init__(
        self,
        exporter: PPTExporter,
        template_path: str = "templates/business_report.pptx",
    ):
        self.exporter = exporter
        self.template_path = template_path

    def create_from_template(
        self,
        title: str = "演示文稿",
        slides_data: Optional[List[Dict[str, Any]]] = None,
    ) -> Dict[str, Any]:
        """使用模板创建PPT并进行内容修改。"""
        prs = self._load_resource()

        # 核心步骤 1: 清理现有幻灯片 (保留母版样式)
        self._clear_existing_slides(prs)

        # 核心步骤 2: 构建内容
        if not slides_data:
            self._create_title_slide(prs, title)
        else:
            for i, slide_info in enumerate(slides_data):
                try:
                    # 容错机制：即使单张失败，也要继续尝试下一张
                    self._create_slide_by_type(prs, slide_info, i)
                except Exception as e:
                    logger.error(f"创建第 {i+1} 张幻灯片失败: {str(e)}")

        return self.exporter.export_ppt(prs, title=title)

    def _load_resource(self) -> Presentation:
        """加载模板资源，带有容错处理。"""
        if os.path.exists(self.template_path):
            logger.info(f"加载模板: {self.template_path}")
            return Presentation(self.template_path)
        else:
            logger.warning(f"模板未找到: {self.template_path}，降级为创建空白PPT")
            return Presentation()

    def _clear_existing_slides(self, prs: Presentation) -> None:
        """清空现有幻灯片 (私有 API 操作，高风险)。"""
        try:
            xml_slides = prs.slides._sldIdLst
            slides_count = len(xml_slides)
            # 倒序删除
            for i in range(slides_count - 1, -1, -1):
                rId = xml_slides[i].rId
                prs.part.drop_rel(rId)
                del xml_slides[i]
        except Exception as e:
            logger.error(f"清理幻灯片时发生错误 (非致命): {str(e)}")

    def _get_layout(self, prs: Presentation, layout_key: str) -> SlideLayout:
        """根据名称安全获取布局 (优先名称匹配，其次索引回退)。"""
        target_name = self.LAYOUT_NAMES.get(layout_key)

        # 1. 尝试通过名称查找 (Plan A)
        for layout in prs.slide_layouts:
            if layout.name == target_name:
                return layout

        # 2. 回退到索引查找 (Plan B)
        fallback_index = self.LAYOUT_INDICES.get(layout_key, 1)
        logger.warning(f"布局 '{target_name}' 未找到，回退使用索引 {fallback_index}")

        if fallback_index < len(prs.slide_layouts):
            return prs.slide_layouts[fallback_index]

        # 3. 最终回退 (保底)
        logger.error("请求的布局索引超出范围，使用默认布局 [0]")
        return prs.slide_layouts[0]

    def _get_body_placeholders(self, slide: Union[Slide, SlideLayout]) -> List[Any]:
        """智能获取正文占位符，过滤掉 Title/Slide Number 等无关框。"""
        body_placeholders = []
        for p in slide.placeholders:
            name = p.name
            # 过滤掉所有标题、图片、页码、表格相关的占位符
            if not any(
                keyword in name
                for keyword in ["Title", "Picture", "Slide Number", "Table"]
            ):
                body_placeholders.append(p)

        # 按位置排序：确保左侧内容填入左边的框 (top, left)
        body_placeholders.sort(key=lambda p: (p.top, p.left))

        # 如果没有找到精确的非标题框，尝试放宽条件 (只排除 Title 和 Slide Number)
        if not body_placeholders:
            for p in slide.placeholders:
                name = p.name
                if not any(keyword in name for keyword in ["Title", "Slide Number"]):
                    body_placeholders.append(p)
            body_placeholders.sort(key=lambda p: (p.top, p.left))

        return body_placeholders

    def _create_title_slide(self, prs: Presentation, title: str, subtitle: str = ""):
        """创建标题幻灯片。"""
        layout = self._get_layout(prs, "title")
        slide = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # 使用动态查找第一个可用内容框作为副标题
        body_placeholders = self._get_body_placeholders(slide)
        if subtitle and body_placeholders:
            self._set_placeholder_content(body_placeholders[0], subtitle)

    def _create_content_slide(self, prs: Presentation, title: str, content: Any):
        """创建内容幻灯片。"""
        layout = self._get_layout(prs, "content")
        slide = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # 使用动态查找第一个可用内容框
        body_placeholders = self._get_body_placeholders(slide)
        if body_placeholders:
            self._set_placeholder_content(body_placeholders[0], content)
        else:
            logger.warning(
                f"幻灯片 '{title}' 的布局 ({layout.name}) 缺少可用内容占位符。"
            )

    def _create_two_column_slide(
        self, prs: Presentation, title: str, left_content: Any, right_content: Any
    ):
        """创建两栏幻灯片。"""
        layout = self._get_layout(prs, "two_column")
        slide = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # 使用动态查找和排序后的占位符
        body_placeholders = self._get_body_placeholders(slide)

        # 智能分配内容
        if len(body_placeholders) >= 1:
            self._set_placeholder_content(body_placeholders[0], left_content)

        if len(body_placeholders) >= 2:
            self._set_placeholder_content(body_placeholders[1], right_content)
        elif len(body_placeholders) == 1:
            logger.warning(f"两栏幻灯片 '{title}' 只有一个内容框，右侧内容将被忽略。")
        else:
            logger.warning(f"两栏幻灯片 '{title}' 没有任何可用内容框。")

    def _set_placeholder_content(self, placeholder: Any, content: Any):
        """设置占位符内容，支持文本和列表。"""
        if not hasattr(placeholder, "text_frame"):
            return

        text_frame = placeholder.text_frame
        text_frame.clear()  # 清空默认文本

        if isinstance(content, list):
            # 将列表项转为带项目符号的段落
            for i, item in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = str(item)
        else:
            text_frame.text = str(content)

    def _create_slide_by_type(
        self, prs: Presentation, slide_info: Dict[str, Any], index: int
    ):
        """工厂分发方法。"""
        slide_type = slide_info.get("slide_type", "content")
        title = slide_info.get("title", f"幻灯片 {index+1}")

        if slide_type == "title":
            self._create_title_slide(prs, title, slide_info.get("subtitle", ""))
        elif slide_type == "two_column":
            self._create_two_column_slide(
                prs,
                title,
                slide_info.get("left_content", ""),
                slide_info.get("right_content", ""),
            )
        else:
            self._create_content_slide(prs, title, slide_info.get("content", ""))


# 使用示例 (Usage Example)
if __name__ == "__main__":
    # 激活日志配置
    logging.basicConfig(
        level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
    )

    import io

    # 1. 实例化依赖 (使用 Mock)
    ppt_exporter = PPTExporter()

    # 2. 注入依赖并创建服务
    # 确保 template_path 指向实际存在的 pptx 文件
    engine = TemplateEngine(
        exporter=ppt_exporter, template_path="templates/business_report.pptx"
    )

    # 3. 准备测试数据 (NVIDIA FY2026 Q2 财报数据)
    sample_data = [
        {
            "slide_type": "title",
            "title": "2025年全球电动汽车(EV)市场分析报告",
            "subtitle": "技术、市场份额与未来趋势",
        },
        {
            "slide_type": "two_column",
            "title": "全球销量领导者 (2025年1-9月)",
            "left_content": [
                "比亚迪 (BYD):",
                "纯电汽车 (BEV) 累计销量: 约 161 万辆。",
                "预计 2025 年全球市场份额将达到 15.7%。",
                "已连续四个季度保持销量领先。",
            ],
            "right_content": [
                "特斯拉 (Tesla):",
                "纯电汽车 (BEV) 累计销量: 约 122 万辆。",
                "同期销量落后比亚迪约 38.8 万辆。",
                "Q3 销量为 497,100 辆，同比增长 7.4%。",
            ],
        },
        {
            "slide_type": "two_column",
            "title": "区域市场动态 (2024-2025)",
            "left_content": [
                "中国:",
                "2024年占全球电动汽车销量的近三分之二。",
                "继续引领全球电动汽车和充电设施的部署。",
            ],
            "right_content": [
                "欧洲:",
                "2024年销量停滞，受主要市场（如德国）补贴减少影响。",
                "2025 年新 CO2 目标生效，预计将推动增长。",
                "挪威仍是领导者，BEV 占新车销量 80% 以上。",
            ],
            "subtitle": "中美欧市场表现分化",
        },
        {
            "slide_type": "two_column",
            "title": "电池技术对比：固态电池 vs 锂离子",
            "left_content": [
                "固态电池 (SSB):",
                "电解质: 固体（陶瓷、聚合物）。",
                "能量密度 (潜力): 350–700 Wh/kg。",
                "安全性: 更高，无热失控风险。",
                "充电速度 (目标): 12-15 分钟至 80%。",
                "材料: 需多 35% 的锂，但显着减少石墨和钴。",
            ],
            "right_content": [
                "传统锂离子 (LIB):",
                "电解质: 液体或凝胶（易燃）。",
                "能量密度 (当前): 150–300 Wh/kg。",
                "安全性: 存在热失控和起火风险。",
                "充电速度 (当前): 30-45 分钟至 80%。",
                "材料: 依赖钴和镍，回收工艺能耗高。",
            ],
        },
        {
            "slide_type": "two_column",
            "title": "关键电池材料供需 (2025)",
            "left_content": [
                "需求增长:",
                "预计 2025-2035 年，关键电池材料市场价值将增长三倍 (CAGR 10.6%)。",
                "镍 (Nickel) 的需求增长迅速，预计将成为继锂和钴之后的主要材料。",
            ],
            "right_content": [
                "供应挑战:",
                "钴 (Cobalt) 需求增长预计最慢，因 LCO 电池市场份额下降。",
                "高品位矿源逐渐枯竭，开采转向低品位矿石，增加了经济和技术挑战。",
                "短期内关键矿物供应迅速增加，存在短期过剩担忧。",
            ],
        },
        {
            "slide_type": "two_column",
            "title": "充电基础设施增长",
            "left_content": [
                "全球概况:",
                "截至 2024 年底，全球公共充电站超 500 万个，两年内翻倍。",
                "中国拥有全球约 65% 的公共充电桩。",
            ],
            "right_content": [
                "区域差异:",
                "美国和英国的 EV 与充电桩比例较高，基础设施部署落后于车辆普及速度。",
                "欧盟法规要求从 2025 年起，主要运输路线每 60 公里必须安装充电点。",
            ],
        },
        {
            "slide_type": "two_column",
            "title": "充电标准与网络联盟 (2025)",
            "left_content": [
                "北美充电标准 (NACS):",
                "特斯拉的标准被广泛采用。",
                "从 2025 年开始，福特、通用、本田等公司的新车型将包含 NACS 端口，以接入特斯拉超级充电网络。",
            ],
            "right_content": [
                "OEM 充电联盟:",
                "七家主要汽车制造商（宝马、通用、本田等）于 2023 年成立合资企业。",
                "计划在北美安装 30,000+ 个大功率充电桩，同时支持 CCS 和 NACS 连接器。",
                "Plug & Charge (ISO 15118) 功能正逐步推广，实现自动认证和计费。",
            ],
        },
        {
            "slide_type": "two_column",
            "title": "全球主要政府政策 (2024-2025)",
            "left_content": [
                "美国:",
                "2024 年初修改清洁车辆税收抵免政策。",
                "允许在销售点即时获得折扣（新车最高 7500 美元）。",
                "但符合条件的 EV 车型减少（约 110 款中的 20 款）。",
            ],
            "right_content": [
                "欧洲:",
                "德国等市场取消或减少补贴导致 2024 年销量停滞。",
                "挪威目标在 2025 年实现 100% 零排放汽车销售。",
                "欧盟 2025 年起实施更严格的 CO2 排放标准。",
            ],
        },
        {
            "slide_type": "two_column",
            "title": "主流车型对比 (VW ID.4 vs Tesla Model Y)",
            "left_content": [
                "大众 ID.4 Pro (2025):",
                "EPA 续航: 约 291 英里。",
                "0-60 mph: 约 5.7 秒 (AWD)。",
                "技术: 支持无线 Apple CarPlay® 和 Android Auto™，保留了部分物理按键。",
                "内部: 空间宽敞，乘坐舒适性高。",
            ],
            "right_content": [
                "特斯拉 Model Y Long Range:",
                "EPA 续航: 约 310 英里。",
                "0-60 mph: 约 4.6 秒。",
                "技术: 完全依赖中央触摸屏，无 CarPlay/Android Auto。",
                "内部: 极简设计，提供紧凑的第三排座位选项。",
            ],
        },
        {
            "slide_type": "two_column",
            "title": "自动驾驶技术路径对比",
            "left_content": [
                "特斯拉 (Tesla FSD):",
                "传感器: 纯视觉（仅 8 个摄像头）。",
                "地图: 不依赖高精地图（HD Maps），设计用于随处行驶。",
                "架构: 端到端 AI，依赖强大的数据车队和 Dojo 训练。",
                "模式: ADAS（高级驾驶辅助），需驾驶员监控。",
            ],
            "right_content": [
                "Waymo:",
                "传感器: 冗余堆栈（29 个摄像头、6 个雷达、5 个 LiDAR）。",
                "地图: 严重依赖高精地图，仅在已绘制的地理围栏区域内运行。",
                "架构: 多个神经网络分离训练，安全性更高。",
                "模式: Robotaxi（自动驾驶出租车），无需驾驶员介入。",
            ],
        },
        {
            "slide_type": "two_column",
            "title": "消费者接受度与主要障碍 (2025)",
            "left_content": [
                "主要障碍 (全球):",
                "成本: 60% 的消费者认为 BEV 仍然过于昂贵。",
                "充电时间: 56% 的消费者认为是主要障碍。",
                "充电站可用性: 54% 的消费者认为是障碍（美国为 44%）。",
            ],
            "right_content": [
                "市场认知:",
                "对成本、维护和租赁的误解减缓了采用速度。",
                "二手 EV 市场迅速扩大，受租约到期和税收抵免推动。",
                "近 30% 的消费者表示愿意等待 30-60 分钟充电。",
            ],
        },
        {
            "slide_type": "content",
            "title": "市场反馈与展望 (至 2030 年)",
            "content": [
                "市场规模预测: 全球电动汽车市场预计将从 2024 年的 8964.5 亿美元增长到 2030 年的 2.82 万亿美元以上。",
                "增长率 (CAGR): 预计 2025 年至 2030 年的复合年增长率将超过 21.51%。",
                "关键驱动力: 持续的技术进步（特别是电池）、政府激励措施和碳排放法规、以及消费者对可持续交通的环保意识增强。",
                "技术趋势: 固态电池等替代技术预计将进一步降低成本并提高性能。",
            ],
        },
    ]

    print("正在生成 PPT (Google Q3 财报)...")

    # 4. 运行完整生成流程
    output_title = "Google_Q3_2025_Report"
    result = engine.create_from_template(title=output_title, slides_data=sample_data)

    # 5. 保存到磁盘
    if result and "buffer" in result:
        ppt_buffer = result["buffer"]
        output_filename = f"{output_title}.pptx"

        try:
            if hasattr(ppt_buffer, "seek"):
                ppt_buffer.seek(0)

            with open(output_filename, "wb") as f:
                f.write(
                    ppt_buffer.getbuffer()
                    if hasattr(ppt_buffer, "getbuffer")
                    else ppt_buffer.read()
                )

            print(f"SUCCESS: 文件已保存至本地 -> {os.path.abspath(output_filename)}")

        except IOError as e:
            print(f"ERROR: 文件写入失败: {e}")
    else:
        print("ERROR: 生成结果无效或缺少文件缓冲区数据")
